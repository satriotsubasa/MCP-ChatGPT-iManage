"""
Document processing module for iManage Deep Research MCP Server
"""

import io
from typing import Optional

# Document processing libraries
try:
    import PyPDF2
    import pdfplumber
    from docx import Document as DocxDocument
    import openpyxl
    from pptx import Presentation
    from bs4 import BeautifulSoup
    DOCUMENT_PROCESSING_AVAILABLE = True
    print("✅ Document processing libraries loaded successfully")
except ImportError as e:
    print(f"⚠️ Some document processing libraries not available: {e}")
    DOCUMENT_PROCESSING_AVAILABLE = False

async def extract_text_from_pdf(content: bytes) -> str:
    """Extract text from PDF content"""
    try:
        print("📄 Attempting PDF text extraction with PyPDF2...")
        
        # Try with PyPDF2 first
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
        text_parts = []
        
        for page_num, page in enumerate(pdf_reader.pages):
            try:
                page_text = page.extract_text()
                if page_text.strip():
                    text_parts.append(f"[Page {page_num + 1}]\n{page_text}")
                    print(f"📄 Extracted text from PDF page {page_num + 1}")
            except Exception as page_error:
                print(f"⚠️ Failed to extract text from PDF page {page_num + 1}: {str(page_error)}")
        
        if text_parts:
            return "\n\n".join(text_parts)
        
        # Fallback to pdfplumber if PyPDF2 didn't work
        print("📄 Fallback: Attempting PDF text extraction with pdfplumber...")
        with pdfplumber.open(io.BytesIO(content)) as pdf:
            text_parts = []
            for page_num, page in enumerate(pdf.pages):
                try:
                    page_text = page.extract_text()
                    if page_text and page_text.strip():
                        text_parts.append(f"[Page {page_num + 1}]\n{page_text}")
                        print(f"📄 Extracted text from PDF page {page_num + 1} with pdfplumber")
                except Exception as page_error:
                    print(f"⚠️ pdfplumber failed on page {page_num + 1}: {str(page_error)}")
            
            return "\n\n".join(text_parts) if text_parts else "PDF content could not be extracted"
            
    except Exception as e:
        print(f"❌ PDF extraction failed: {str(e)}")
        return f"PDF processing error: {str(e)}"

async def extract_text_from_docx(content: bytes) -> str:
    """Extract text from Word document"""
    try:
        print("📄 Attempting Word document text extraction...")
        
        doc = DocxDocument(io.BytesIO(content))
        text_parts = []
        
        # Extract paragraph text
        for para_num, paragraph in enumerate(doc.paragraphs):
            if paragraph.text.strip():
                text_parts.append(paragraph.text)
        
        # Extract table text
        for table_num, table in enumerate(doc.tables):
            print(f"📊 Processing table {table_num + 1}")
            table_text = []
            for row in table.rows:
                row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                if row_text:
                    table_text.append(row_text)
            
            if table_text:
                text_parts.append(f"\n[Table {table_num + 1}]\n" + "\n".join(table_text))
        
        extracted_text = "\n".join(text_parts)
        print(f"✅ Extracted {len(extracted_text)} characters from Word document")
        return extracted_text if extracted_text.strip() else "No readable text found in Word document"
        
    except Exception as e:
        print(f"❌ Word document extraction failed: {str(e)}")
        return f"Word document processing error: {str(e)}"

async def extract_text_from_excel(content: bytes) -> str:
    """Extract text from Excel spreadsheet"""
    try:
        print("📊 Attempting Excel spreadsheet text extraction...")
        
        workbook = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
        text_parts = []
        
        for sheet_name in workbook.sheetnames:
            print(f"📊 Processing Excel sheet: {sheet_name}")
            sheet = workbook[sheet_name]
            
            sheet_data = []
            for row_num, row in enumerate(sheet.iter_rows(values_only=True), 1):
                row_data = [str(cell) if cell is not None else "" for cell in row]
                if any(cell.strip() for cell in row_data):  # Skip empty rows
                    sheet_data.append(" | ".join(row_data))
                
                # Limit to first 100 rows per sheet to avoid huge outputs
                if row_num > 100:
                    sheet_data.append("... (truncated, showing first 100 rows)")
                    break
            
            if sheet_data:
                text_parts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(sheet_data))
        
        extracted_text = "\n\n".join(text_parts)
        print(f"✅ Extracted {len(extracted_text)} characters from Excel file")
        return extracted_text if extracted_text.strip() else "No readable data found in Excel file"
        
    except Exception as e:
        print(f"❌ Excel extraction failed: {str(e)}")
        return f"Excel processing error: {str(e)}"

async def extract_text_from_pptx(content: bytes) -> str:
    """Extract text from PowerPoint presentation"""
    try:
        print("📊 Attempting PowerPoint text extraction...")
        
        prs = Presentation(io.BytesIO(content))
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if slide_text:
                text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text))
                print(f"📊 Extracted text from PowerPoint slide {slide_num}")
        
        extracted_text = "\n\n".join(text_parts)
        print(f"✅ Extracted {len(extracted_text)} characters from PowerPoint")
        return extracted_text if extracted_text.strip() else "No readable text found in PowerPoint"
        
    except Exception as e:
        print(f"❌ PowerPoint extraction failed: {str(e)}")
        return f"PowerPoint processing error: {str(e)}"

async def extract_text_from_html(content: bytes) -> str:
    """Extract text from HTML content"""
    try:
        print("🌐 Attempting HTML text extraction...")
        
        soup = BeautifulSoup(content, 'html.parser')
        
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()
        
        # Get text
        text = soup.get_text()
        
        # Clean up whitespace
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = '\n'.join(chunk for chunk in chunks if chunk)
        
        print(f"✅ Extracted {len(text)} characters from HTML")
        return text if text.strip() else "No readable text found in HTML"
        
    except Exception as e:
        print(f"❌ HTML extraction failed: {str(e)}")
        return f"HTML processing error: {str(e)}"

async def process_document_content(content: bytes, content_type: str, filename: str = "") -> str:
    """Process document content based on content type and extract readable text"""
    
    content_type_lower = content_type.lower()
    filename_lower = filename.lower()
    
    print(f"📋 Processing document: type={content_type}, filename={filename}, size={len(content)} bytes")
    
    # Determine document type and process accordingly
    if "pdf" in content_type_lower or filename_lower.endswith('.pdf'):
        return await extract_text_from_pdf(content)
    
    elif "msword" in content_type_lower or "officedocument.wordprocessingml" in content_type_lower or filename_lower.endswith(('.doc', '.docx')):
        return await extract_text_from_docx(content)
    
    elif "excel" in content_type_lower or "spreadsheetml" in content_type_lower or filename_lower.endswith(('.xls', '.xlsx')):
        return await extract_text_from_excel(content)
    
    elif "powerpoint" in content_type_lower or "presentationml" in content_type_lower or filename_lower.endswith(('.ppt', '.pptx')):
        return await extract_text_from_pptx(content)
    
    elif "html" in content_type_lower or filename_lower.endswith(('.html', '.htm')):
        return await extract_text_from_html(content)
    
    elif "text" in content_type_lower or "json" in content_type_lower or "xml" in content_type_lower:
        try:
            text_content = content.decode('utf-8', errors='ignore')
            print(f"✅ Decoded text content: {len(text_content)} characters")
            return text_content
        except Exception as e:
            print(f"❌ Text decoding failed: {str(e)}")
            return f"Text decoding error: {str(e)}"
    
    else:
        print(f"⚠️ Unsupported document type: {content_type}")
        return f"Unsupported document format: {content_type}. File size: {len(content)} bytes. Unable to extract text content."

def get_processing_capabilities():
    """Get current document processing capabilities"""
    capabilities = {
        "libraries_available": {
            "PyPDF2": False,
            "pdfplumber": False,
            "python-docx": False,
            "openpyxl": False,
            "python-pptx": False,
            "beautifulsoup4": False
        },
        "supported_formats": [],
        "overall_status": DOCUMENT_PROCESSING_AVAILABLE
    }
    
    # Test each library
    try:
        import PyPDF2
        capabilities["libraries_available"]["PyPDF2"] = True
        capabilities["supported_formats"].append("PDF")
    except ImportError:
        pass
    
    try:
        import pdfplumber
        capabilities["libraries_available"]["pdfplumber"] = True
    except ImportError:
        pass
    
    try:
        from docx import Document
        capabilities["libraries_available"]["python-docx"] = True
        capabilities["supported_formats"].append("Word Documents (.docx)")
    except ImportError:
        pass
    
    try:
        import openpyxl
        capabilities["libraries_available"]["openpyxl"] = True
        capabilities["supported_formats"].append("Excel Spreadsheets (.xlsx)")
    except ImportError:
        pass
    
    try:
        from pptx import Presentation
        capabilities["libraries_available"]["python-pptx"] = True
        capabilities["supported_formats"].append("PowerPoint Presentations (.pptx)")
    except ImportError:
        pass
    
    try:
        from bs4 import BeautifulSoup
        capabilities["libraries_available"]["beautifulsoup4"] = True
        capabilities["supported_formats"].append("HTML Documents")
    except ImportError:
        pass
    
    capabilities["supported_formats"].extend(["Plain Text", "JSON", "XML"])
    
    return capabilities