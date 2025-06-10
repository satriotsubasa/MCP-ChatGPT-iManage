#!/usr/bin/env python3
"""
iManage Deep Research MCP Server for ChatGPT Integration

This MCP server enables ChatGPT to perform deep research using iManage Work API.
It provides search and fetch capabilities for documents stored in iManage.

Features:
- Title-based search
- Keyword search (full-text)
- Document retrieval and download
- Token-based authentication with caching
- Comprehensive logging with emojis
- Support for multiple search strategies

Environment Variables Required:
- AUTH_URL_PREFIX: iManage authentication URL prefix
- URL_PREFIX: iManage API URL prefix 
- _USERNAME: iManage username
- PASSWORD: iManage password
- CLIENT_ID: OAuth client ID
- CLIENT_SECRET: OAuth client secret
- CUSTOMER_ID: iManage customer ID
- LIBRARY_ID: iManage library ID (default library to search)
"""

import asyncio
import httpx
import json
import os
import time
import io
import tempfile
from typing import Dict, List, Any, Optional
from fastapi import FastAPI, HTTPException, Request
from fastapi.responses import StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, ValidationError
import logging

# Document processing libraries
try:
    import PyPDF2
    import pdfplumber
    from docx import Document as DocxDocument
    import openpyxl
    from pptx import Presentation
    from bs4 import BeautifulSoup
    DOCUMENT_PROCESSING_AVAILABLE = True
    print("✅ Document processing libraries loaded successfully")
except ImportError as e:
    print(f"⚠️ Some document processing libraries not available: {e}")
    DOCUMENT_PROCESSING_AVAILABLE = False

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = FastAPI(title="iManage Deep Research MCP Server")

# Add CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ---- Configuration ----
AUTH_URL_PREFIX = os.getenv("AUTH_URL_PREFIX", "")
URL_PREFIX = os.getenv("URL_PREFIX", "")
_USERNAME = os.getenv("_USERNAME", "")
PASSWORD = os.getenv("PASSWORD", "")
CLIENT_ID = os.getenv("CLIENT_ID", "")
CLIENT_SECRET = os.getenv("CLIENT_SECRET", "")
CUSTOMER_ID = os.getenv("CUSTOMER_ID", "")
LIBRARY_ID = os.getenv("LIBRARY_ID", "")

# Validate required environment variables
required_vars = [
    "AUTH_URL_PREFIX", "URL_PREFIX", "_USERNAME", "PASSWORD", 
    "CLIENT_ID", "CLIENT_SECRET", "CUSTOMER_ID", "LIBRARY_ID"
]

for var in required_vars:
    if not globals()[var]:
        raise ValueError(f"❌ Required environment variable {var} is not set")

print("✅ All required environment variables are configured")

# ---- Token cache ----
token_cache = {"token": None, "expires": 0}

async def get_token() -> str:
    """Get authentication token with caching"""
    if token_cache["token"] and token_cache["expires"] > time.time():
        print("🔓 Using cached access token")
        return token_cache["token"]
    
    print("🔐 Authenticating to iManage...")
    auth_url = f"{AUTH_URL_PREFIX}/oauth2/token?scope=admin"
    data = {
        "username": _USERNAME,
        "password": PASSWORD,
        "grant_type": "password",
        "client_id": CLIENT_ID,
        "client_secret": CLIENT_SECRET
    }
    headers = {"Content-Type": "application/x-www-form-urlencoded"}
    
    try:
        async with httpx.AsyncClient(timeout=30.0) as client:
            res = await client.post(auth_url, data=data, headers=headers)
            res.raise_for_status()
            token_data = res.json()
            token_cache["token"] = token_data["access_token"]
            token_cache["expires"] = time.time() + token_data.get("expires_in", 1800) - 60
            print("✅ Authentication successful")
            return token_data["access_token"]
    except Exception as e:
        print(f"❌ Authentication failed: {str(e)}")
        raise HTTPException(status_code=401, detail=f"Authentication failed: {str(e)}")

# ---- Simple Models ----
class SearchResult(BaseModel):
    id: str
    title: str
    text: str
    url: Optional[str] = None
    metadata: Optional[Dict[str, str]] = None

# ---- Core Search Functions ----
async def search_documents_title(query: str, limit: int = 20) -> List[SearchResult]:
    """Search documents by title/name"""
    print(f"🔍 Searching documents by title: '{query}'")
    
    token = await get_token()
    search_url = f"{URL_PREFIX}/api/v2/customers/{CUSTOMER_ID}/libraries/{LIBRARY_ID}/documents/search"
    
    headers = {
        "X-Auth-Token": token,
        "Content-Type": "application/json"
    }
    
    # Search using POST with filters for title search
    search_body = {
        "limit": limit,
        "filters": {
            "name": query
        },
        "profile_fields": {
            "document": [
                "id", "name", "document_number", "version", "author", 
                "edit_date", "create_date", "size", "type"
            ]
        }
    }
    
    try:
        async with httpx.AsyncClient(timeout=60.0) as client:
            print(f"🔍 Sending title search request: {json.dumps(search_body, indent=2)}")
            response = await client.post(search_url, headers=headers, json=search_body)
            
            # Log the response for debugging
            print(f"📊 Response status: {response.status_code}")
            print(f"📊 Response headers: {dict(response.headers)}")
            
            if response.status_code == 400:
                error_text = response.text
                print(f"❌ 400 Error details: {error_text}")
                # Try a simpler search format
                return await search_documents_simple(query, limit, "title")
            
            response.raise_for_status()
            data = response.json()
            
            results = []
            for doc in data.get("data", []):
                doc_id = doc.get("id", "")
                title = doc.get("name", "Untitled Document")
                
                # Create text snippet from document metadata
                text_parts = []
                if doc.get("comments"):
                    text_parts.append(f"Comments: {doc['comments']}")
                if doc.get("author"):
                    text_parts.append(f"Author: {doc['author']}")
                if doc.get("type"):
                    text_parts.append(f"Type: {doc['type']}")
                if doc.get("edit_date"):
                    text_parts.append(f"Last Modified: {doc['edit_date']}")
                
                text = "; ".join(text_parts) if text_parts else "Document metadata"
                
                # Generate document URL for citations
                doc_url = f"{URL_PREFIX}/api/v2/customers/{CUSTOMER_ID}/libraries/{LIBRARY_ID}/documents/{doc_id}"
                
                metadata = {
                    "document_number": str(doc.get("document_number", "")),
                    "version": str(doc.get("version", "")),
                    "size": str(doc.get("size", "")),
                    "search_type": "title"
                }
                
                results.append(SearchResult(
                    id=doc_id,
                    title=title,
                    text=text,
                    url=doc_url,
                    metadata=metadata
                ))
            
            print(f"📄 Found {len(results)} documents by title")
            return results
            
    except Exception as e:
        print(f"❌ Title search failed: {str(e)}")
        # Try fallback search methods
        try:
            return await search_documents_simple(query, limit, "title")
        except Exception as fallback_error:
            print(f"❌ Fallback search also failed: {str(fallback_error)}")
            return []

async def search_documents_simple(query: str, limit: int = 20, search_type: str = "simple") -> List[SearchResult]:
    """Simple search using GET parameters - fallback method"""
    print(f"🔍 Trying simple search ({search_type}): '{query}'")
    
    token = await get_token()
    
    # Use GET endpoint with query parameters
    search_url = f"{URL_PREFIX}/api/v2/customers/{CUSTOMER_ID}/libraries/{LIBRARY_ID}/documents"
    
    headers = {"X-Auth-Token": token}
    
    # Try different parameter combinations
    params_options = [
        {"name": query, "limit": limit},
        {"anywhere": query, "limit": limit},
        {"title": query, "limit": limit},
        {"body": query, "limit": limit},
        {"q": query, "limit": limit}  # Sometimes 'q' is used as generic search
    ]
    
    for params in params_options:
        try:
            print(f"🔍 Trying search with params: {params}")
            async with httpx.AsyncClient(timeout=60.0) as client:
                response = await client.get(search_url, headers=headers, params=params)
                
                if response.status_code == 200:
                    try:
                        data = response.json()
                    except json.JSONDecodeError:
                        print(f"⚠️ Response is not JSON: {response.text[:200]}")
                        continue
                    
                    # Handle different response formats
                    documents = []
                    if isinstance(data, dict):
                        documents = data.get("data", [])
                    elif isinstance(data, list):
                        documents = data
                    
                    results = []
                    for doc in documents:
                        if not isinstance(doc, dict):
                            continue
                            
                        doc_id = doc.get("id", "")
                        title = doc.get("name", "Untitled Document")
                        
                        # Create text snippet from document metadata
                        text_parts = []
                        if doc.get("comments"):
                            text_parts.append(f"Comments: {doc['comments']}")
                        if doc.get("author"):
                            text_parts.append(f"Author: {doc['author']}")
                        if doc.get("type"):
                            text_parts.append(f"Type: {doc['type']}")
                        if doc.get("edit_date"):
                            text_parts.append(f"Last Modified: {doc['edit_date']}")
                        
                        text = "; ".join(text_parts) if text_parts else f"Document found with {search_type} search"
                        
                        # Generate document URL for citations - FIXED URL FORMAT
                        doc_url = f"{URL_PREFIX}/work/web/api/v2/customers/{CUSTOMER_ID}/libraries/{LIBRARY_ID}/documents/{doc_id}"
                        
                        metadata = {
                            "document_number": str(doc.get("document_number", "")),
                            "version": str(doc.get("version", "")),
                            "size": str(doc.get("size", "")),
                            "search_type": search_type
                        }
                        
                        results.append(SearchResult(
                            id=doc_id,
                            title=title,
                            text=text,
                            url=doc_url,
                            metadata=metadata
                        ))
                    
                    print(f"📄 Simple search found {len(results)} documents")
                    return results[:limit]  # Limit results
                else:
                    print(f"⚠️ Search params {params} returned {response.status_code}: {response.text[:200]}")
                    
        except Exception as e:
            print(f"⚠️ Search params {params} failed: {str(e)}")
            continue
    
    print("❌ All simple search methods failed")
    return []
async def search_documents_keyword(query: str, limit: int = 20) -> List[SearchResult]:
    """Search documents by keywords (full-text search)"""
    print(f"🔍 Searching documents by keywords: '{query}'")
    
    token = await get_token()
    search_url = f"{URL_PREFIX}/api/v2/customers/{CUSTOMER_ID}/libraries/{LIBRARY_ID}/documents/search"
    
    headers = {
        "X-Auth-Token": token,
        "Content-Type": "application/json"
    }
    
    # Try different search body formats
    search_body_options = [
        # Option 1: Basic fields only (removing comments)
        {
            "limit": limit,
            "filters": {
                "anywhere": query
            },
            "profile_fields": {
                "document": [
                    "id", "name", "document_number", "version", "author", 
                    "edit_date", "create_date", "size", "type"
                ]
            }
        },
        # Option 2: Even more basic fields
        {
            "limit": limit,
            "filters": {
                "anywhere": query
            },
            "profile_fields": {
                "document": [
                    "id", "name", "author", "edit_date", "type"
                ]
            }
        },
        # Option 3: No profile_fields specified
        {
            "limit": limit,
            "filters": {
                "anywhere": query
            }
        },
        # Option 4: Body search only
        {
            "limit": limit,
            "filters": {
                "body": query
            }
        }
    ]
    
    for i, search_body in enumerate(search_body_options):
        try:
            print(f"🔍 Trying keyword search format {i+1}: {json.dumps(search_body, indent=2)}")
            async with httpx.AsyncClient(timeout=60.0) as client:
                response = await client.post(search_url, headers=headers, json=search_body)
                
                if response.status_code == 200:
                    data = response.json()
                    
                    results = []
                    for doc in data.get("data", []):
                        doc_id = doc.get("id", "")
                        title = doc.get("name", "Untitled Document")
                        
                        # Create text snippet from document metadata
                        text_parts = []
                        if doc.get("comments"):
                            text_parts.append(f"Comments: {doc['comments']}")
                        if doc.get("author"):
                            text_parts.append(f"Author: {doc['author']}")
                        if doc.get("type"):
                            text_parts.append(f"Type: {doc['type']}")
                        if doc.get("edit_date"):
                            text_parts.append(f"Last Modified: {doc['edit_date']}")
                        
                        text = "; ".join(text_parts) if text_parts else f"Document contains keyword: {query}"
                        
                        # Generate document URL for citations - FIXED URL FORMAT
                        doc_url = f"{URL_PREFIX}/work/web/api/v2/customers/{CUSTOMER_ID}/libraries/{LIBRARY_ID}/documents/{doc_id}"
                        
                        metadata = {
                            "document_number": str(doc.get("document_number", "")),
                            "version": str(doc.get("version", "")),
                            "size": str(doc.get("size", "")),
                            "search_type": "keyword",
                            "search_query": query
                        }
                        
                        results.append(SearchResult(
                            id=doc_id,
                            title=title,
                            text=text,
                            url=doc_url,
                            metadata=metadata
                        ))
                    
                    print(f"📄 Found {len(results)} documents by keywords")
                    return results
                else:
                    print(f"⚠️ Search format {i+1} returned {response.status_code}: {response.text}")
                    
        except Exception as e:
            print(f"⚠️ Search format {i+1} failed: {str(e)}")
            continue
    
    # If all POST methods fail, try GET fallback
    print("🔄 POST search failed, trying GET fallback")
    try:
        return await search_documents_simple(query, limit, "keyword")
    except Exception as fallback_error:
        print(f"❌ Keyword search and fallback failed: {str(fallback_error)}")
        return []

# ---- Document Processing Functions ----
async def extract_text_from_pdf(content: bytes) -> str:
    """Extract text from PDF content"""
    try:
        print("📄 Attempting PDF text extraction with PyPDF2...")
        
        # Try with PyPDF2 first
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
        text_parts = []
        
        for page_num, page in enumerate(pdf_reader.pages):
            try:
                page_text = page.extract_text()
                if page_text.strip():
                    text_parts.append(f"[Page {page_num + 1}]\n{page_text}")
                    print(f"📄 Extracted text from PDF page {page_num + 1}")
            except Exception as page_error:
                print(f"⚠️ Failed to extract text from PDF page {page_num + 1}: {str(page_error)}")
        
        if text_parts:
            return "\n\n".join(text_parts)
        
        # Fallback to pdfplumber if PyPDF2 didn't work
        print("📄 Fallback: Attempting PDF text extraction with pdfplumber...")
        with pdfplumber.open(io.BytesIO(content)) as pdf:
            text_parts = []
            for page_num, page in enumerate(pdf.pages):
                try:
                    page_text = page.extract_text()
                    if page_text and page_text.strip():
                        text_parts.append(f"[Page {page_num + 1}]\n{page_text}")
                        print(f"📄 Extracted text from PDF page {page_num + 1} with pdfplumber")
                except Exception as page_error:
                    print(f"⚠️ pdfplumber failed on page {page_num + 1}: {str(page_error)}")
            
            return "\n\n".join(text_parts) if text_parts else "PDF content could not be extracted"
            
    except Exception as e:
        print(f"❌ PDF extraction failed: {str(e)}")
        return f"PDF processing error: {str(e)}"

async def extract_text_from_docx(content: bytes) -> str:
    """Extract text from Word document"""
    try:
        print("📄 Attempting Word document text extraction...")
        
        doc = DocxDocument(io.BytesIO(content))
        text_parts = []
        
        # Extract paragraph text
        for para_num, paragraph in enumerate(doc.paragraphs):
            if paragraph.text.strip():
                text_parts.append(paragraph.text)
        
        # Extract table text
        for table_num, table in enumerate(doc.tables):
            print(f"📊 Processing table {table_num + 1}")
            table_text = []
            for row in table.rows:
                row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                if row_text:
                    table_text.append(row_text)
            
            if table_text:
                text_parts.append(f"\n[Table {table_num + 1}]\n" + "\n".join(table_text))
        
        extracted_text = "\n".join(text_parts)
        print(f"✅ Extracted {len(extracted_text)} characters from Word document")
        return extracted_text if extracted_text.strip() else "No readable text found in Word document"
        
    except Exception as e:
        print(f"❌ Word document extraction failed: {str(e)}")
        return f"Word document processing error: {str(e)}"

async def extract_text_from_excel(content: bytes) -> str:
    """Extract text from Excel spreadsheet"""
    try:
        print("📊 Attempting Excel spreadsheet text extraction...")
        
        workbook = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
        text_parts = []
        
        for sheet_name in workbook.sheetnames:
            print(f"📊 Processing Excel sheet: {sheet_name}")
            sheet = workbook[sheet_name]
            
            sheet_data = []
            for row_num, row in enumerate(sheet.iter_rows(values_only=True), 1):
                row_data = [str(cell) if cell is not None else "" for cell in row]
                if any(cell.strip() for cell in row_data):  # Skip empty rows
                    sheet_data.append(" | ".join(row_data))
                
                # Limit to first 100 rows per sheet to avoid huge outputs
                if row_num > 100:
                    sheet_data.append("... (truncated, showing first 100 rows)")
                    break
            
            if sheet_data:
                text_parts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(sheet_data))
        
        extracted_text = "\n\n".join(text_parts)
        print(f"✅ Extracted {len(extracted_text)} characters from Excel file")
        return extracted_text if extracted_text.strip() else "No readable data found in Excel file"
        
    except Exception as e:
        print(f"❌ Excel extraction failed: {str(e)}")
        return f"Excel processing error: {str(e)}"

async def extract_text_from_pptx(content: bytes) -> str:
    """Extract text from PowerPoint presentation"""
    try:
        print("📊 Attempting PowerPoint text extraction...")
        
        prs = Presentation(io.BytesIO(content))
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if slide_text:
                text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text))
                print(f"📊 Extracted text from PowerPoint slide {slide_num}")
        
        extracted_text = "\n\n".join(text_parts)
        print(f"✅ Extracted {len(extracted_text)} characters from PowerPoint")
        return extracted_text if extracted_text.strip() else "No readable text found in PowerPoint"
        
    except Exception as e:
        print(f"❌ PowerPoint extraction failed: {str(e)}")
        return f"PowerPoint processing error: {str(e)}"

async def extract_text_from_html(content: bytes) -> str:
    """Extract text from HTML content"""
    try:
        print("🌐 Attempting HTML text extraction...")
        
        soup = BeautifulSoup(content, 'html.parser')
        
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()
        
        # Get text
        text = soup.get_text()
        
        # Clean up whitespace
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = '\n'.join(chunk for chunk in chunks if chunk)
        
        print(f"✅ Extracted {len(text)} characters from HTML")
        return text if text.strip() else "No readable text found in HTML"
        
    except Exception as e:
        print(f"❌ HTML extraction failed: {str(e)}")
        return f"HTML processing error: {str(e)}"

async def process_document_content(content: bytes, content_type: str, filename: str = "") -> str:
    """Process document content based on content type and extract readable text"""
    
    content_type_lower = content_type.lower()
    filename_lower = filename.lower()
    
    print(f"📋 Processing document: type={content_type}, filename={filename}, size={len(content)} bytes")
    
    # Determine document type and process accordingly
    if "pdf" in content_type_lower or filename_lower.endswith('.pdf'):
        return await extract_text_from_pdf(content)
    
    elif "msword" in content_type_lower or "officedocument.wordprocessingml" in content_type_lower or filename_lower.endswith(('.doc', '.docx')):
        return await extract_text_from_docx(content)
    
    elif "excel" in content_type_lower or "spreadsheetml" in content_type_lower or filename_lower.endswith(('.xls', '.xlsx')):
        return await extract_text_from_excel(content)
    
    elif "powerpoint" in content_type_lower or "presentationml" in content_type_lower or filename_lower.endswith(('.ppt', '.pptx')):
        return await extract_text_from_pptx(content)
    
    elif "html" in content_type_lower or filename_lower.endswith(('.html', '.htm')):
        return await extract_text_from_html(content)
    
    elif "text" in content_type_lower or "json" in content_type_lower or "xml" in content_type_lower:
        try:
            text_content = content.decode('utf-8', errors='ignore')
            print(f"✅ Decoded text content: {len(text_content)} characters")
            return text_content
        except Exception as e:
            print(f"❌ Text decoding failed: {str(e)}")
            return f"Text decoding error: {str(e)}"
    
    else:
        print(f"⚠️ Unsupported document type: {content_type}")
        return f"Unsupported document format: {content_type}. File size: {len(content)} bytes. Unable to extract text content."
async def fetch_document_content(doc_id: str) -> Dict[str, Any]:
    """Fetch full document content and metadata with enhanced text extraction"""
    print(f"📥 Fetching document content: {doc_id}")
    
    token = await get_token()
    
    # First get document metadata - FIXED URL FORMAT
    doc_url = f"{URL_PREFIX}/work/web/api/v2/customers/{CUSTOMER_ID}/libraries/{LIBRARY_ID}/documents/{doc_id}"
    headers = {"X-Auth-Token": token}
    
    try:
        async with httpx.AsyncClient(timeout=60.0) as client:
            # Get document metadata
            print(f"📋 Getting document metadata from: {doc_url}")
            response = await client.get(doc_url, headers=headers)
            response.raise_for_status()
            doc_data = response.json().get("data", {})
            
            title = doc_data.get("name", "Untitled Document")
            doc_type = doc_data.get("type", "Unknown")
            doc_size = doc_data.get("size", 0)
            
            print(f"📄 Document metadata: title='{title}', type='{doc_type}', size={doc_size}")
            
            # Try to download document content - FIXED URL FORMAT
            download_url = f"{URL_PREFIX}/work/web/api/v2/customers/{CUSTOMER_ID}/libraries/{LIBRARY_ID}/documents/{doc_id}/download"
            
            document_text = ""
            download_success = False
            
            try:
                print(f"⬇️ Downloading document from: {download_url}")
                download_response = await client.get(download_url, headers=headers)
                download_response.raise_for_status()
                
                content_type = download_response.headers.get("content-type", "").lower()
                content_disposition = download_response.headers.get("content-disposition", "")
                filename = title  # Use document title as filename fallback
                
                # Try to extract filename from content-disposition header
                if "filename=" in content_disposition:
                    try:
                        filename = content_disposition.split("filename=")[1].strip('"')
                    except:
                        pass
                
                print(f"📄 Downloaded: content_type='{content_type}', filename='{filename}', size={len(download_response.content)}")
                
                if DOCUMENT_PROCESSING_AVAILABLE:
                    # Use enhanced document processing
                    document_text = await process_document_content(
                        download_response.content, 
                        content_type, 
                        filename
                    )
                    download_success = True
                    print(f"✅ Successfully processed document: {len(document_text)} characters extracted")
                else:
                    # Fallback to simple text extraction
                    if "text" in content_type or "json" in content_type or "xml" in content_type:
                        document_text = download_response.text
                        download_success = True
                        print(f"✅ Text content extracted: {len(document_text)} characters")
                    else:
                        document_text = f"Binary document ({content_type}). Size: {len(download_response.content)} bytes. Document processing libraries not available for text extraction."
                        print(f"⚠️ Binary document, no processing available")
                
            except Exception as download_error:
                print(f"❌ Document download failed: {str(download_error)}")
                document_text = f"Document download failed: {str(download_error)}. Document metadata available below."
            
            # Build comprehensive document information
            text_parts = []
            
            # Add document content
            if document_text and document_text.strip():
                text_parts.append("=== DOCUMENT CONTENT ===")
                text_parts.append(document_text)
            else:
                text_parts.append("=== DOCUMENT CONTENT UNAVAILABLE ===")
                text_parts.append("Document content could not be extracted or is empty.")
            
            # Add document metadata
            text_parts.append("\n=== DOCUMENT METADATA ===")
            if doc_data.get("comments"):
                text_parts.append(f"Comments: {doc_data['comments']}")
            if doc_data.get("author"):
                text_parts.append(f"Author: {doc_data['author']}")
            if doc_data.get("type"):
                text_parts.append(f"Document Type: {doc_data['type']}")
            if doc_data.get("size"):
                text_parts.append(f"Size: {doc_data['size']} bytes")
            if doc_data.get("edit_date"):
                text_parts.append(f"Last Modified: {doc_data['edit_date']}")
            if doc_data.get("create_date"):
                text_parts.append(f"Created: {doc_data['create_date']}")
            if doc_data.get("document_number"):
                text_parts.append(f"Document Number: {doc_data['document_number']}")
            if doc_data.get("version"):
                text_parts.append(f"Version: {doc_data['version']}")
            
            # Add processing status
            text_parts.append(f"\n=== PROCESSING STATUS ===")
            text_parts.append(f"Download Successful: {download_success}")
            text_parts.append(f"Text Extraction: {'Successful' if document_text and len(document_text) > 100 else 'Limited or Failed'}")
            text_parts.append(f"Document Processing Libraries: {'Available' if DOCUMENT_PROCESSING_AVAILABLE else 'Not Available'}")
            
            full_text = "\n".join(text_parts)
            
            # Generate document URL for citations - FIXED URL FORMAT
            doc_citation_url = f"{URL_PREFIX}/work/web/api/v2/customers/{CUSTOMER_ID}/libraries/{LIBRARY_ID}/documents/{doc_id}"
            
            metadata = {
                "document_number": str(doc_data.get("document_number", "")),
                "version": str(doc_data.get("version", "")),
                "author": doc_data.get("author", ""),
                "type": doc_data.get("type", ""),
                "size": str(doc_data.get("size", "")),
                "download_url": download_url,
                "download_success": str(download_success),
                "text_extracted": str(len(document_text) > 100),
                "processing_available": str(DOCUMENT_PROCESSING_AVAILABLE)
            }
            
            print(f"📊 Document processing complete: {len(full_text)} total characters")
            
            return {
                "id": doc_id,
                "title": title,
                "text": full_text,
                "url": doc_citation_url,
                "metadata": metadata
            }
            
    except Exception as e:
        print(f"❌ Failed to fetch document {doc_id}: {str(e)}")
        error_msg = f"Failed to fetch document {doc_id}: {str(e)}"
        return {
            "id": doc_id,
            "title": f"Error accessing document {doc_id}",
            "text": error_msg,
            "url": f"{URL_PREFIX}/work/web/api/v2/customers/{CUSTOMER_ID}/libraries/{LIBRARY_ID}/documents/{doc_id}",
            "metadata": {"error": str(e)}
        }

# ---- MCP Protocol Handler ----
@app.post("/")
async def mcp_handler(request: Request):
    """Main MCP protocol handler - handles raw JSON"""
    print("📨 MCP request received")
    
    try:
        # Parse the raw JSON request
        body = await request.json()
        print(f"🔍 Request body: {json.dumps(body, indent=2)}")
        
        method = body.get("method", "")
        request_id = body.get("id")
        params = body.get("params", {})
        
        if method == "initialize":
            print("🚀 Initialize request")
            return {
                "jsonrpc": "2.0",
                "id": request_id,
                "result": {
                    "protocolVersion": "2024-11-05",
                    "capabilities": {
                        "tools": {}
                    },
                    "serverInfo": {
                        "name": "iManage Deep Research MCP Server",
                        "version": "1.0.0"
                    },
                    "instructions": "This server provides access to iManage document search and retrieval. No authentication is required as the server handles iManage authentication internally."
                }
            }
        
        elif method == "auth/list":
            print("🔓 Auth methods requested")
            return {
                "jsonrpc": "2.0",
                "id": request_id,
                "result": {
                    "authMethods": []  # Empty array means no auth required
                }
            }
        
        elif method == "auth/status":
            print("🔓 Auth status requested")
            return {
                "jsonrpc": "2.0",
                "id": request_id,
                "result": {
                    "authenticated": True,
                    "method": "none"
                }
            }
        
        elif method == "tools/list":
            print("🛠️ Tools list requested")
            return {
                "jsonrpc": "2.0",
                "id": request_id,
                "result": {
                    "tools": [
                        {
                            "name": "search",
                            "description": "Search for documents in iManage using title search or keyword search. "
                                          "For title search, use specific document names or titles. "
                                          "For keyword search, use terms that might appear in document content. "
                                          "The system will automatically determine the best search strategy. "
                                          "You can search for legal documents, contracts, memos, emails, and other business documents. "
                                          "Use specific terms like client names, matter names, document types, or legal concepts.",
                            "inputSchema": {
                                "type": "object",
                                "properties": {
                                    "query": {
                                        "type": "string", 
                                        "description": "Search query. Can be document titles, keywords, or phrases to search for in documents."
                                    }
                                },
                                "required": ["query"]
                            }
                        },
                        {
                            "name": "fetch",
                            "description": "Retrieve the complete content and metadata of a specific document by its ID. "
                                          "Use this after finding documents with search to get the full document content for analysis.",
                            "inputSchema": {
                                "type": "object",
                                "properties": {
                                    "id": {"type": "string", "description": "Document ID obtained from search results"}
                                },
                                "required": ["id"]
                            }
                        }
                    ]
                }
            }
        
        elif method == "tools/call":
            tool_name = params.get("name")
            arguments = params.get("arguments", {})
            
            print(f"🔧 Tool call: {tool_name} with args: {arguments}")
            
            if tool_name == "search":
                try:
                    query = arguments.get("query", "")
                    if not query:
                        raise ValueError("Query parameter is required")
                    
                    print(f"🔍 Searching for: '{query}'")
                    
                    # Try searches with better error handling
                    all_results = {}
                    
                    # Try title search first
                    try:
                        title_results = await search_documents_title(query, limit=10)
                        for result in title_results:
                            all_results[result.id] = result
                        print(f"✅ Title search returned {len(title_results)} results")
                    except Exception as title_error:
                        print(f"⚠️ Title search failed: {str(title_error)}")
                    
                    # Try keyword search
                    try:
                        keyword_results = await search_documents_keyword(query, limit=10)
                        for result in keyword_results:
                            if result.id not in all_results:
                                all_results[result.id] = result
                        print(f"✅ Keyword search returned {len(keyword_results)} results")
                    except Exception as keyword_error:
                        print(f"⚠️ Keyword search failed: {str(keyword_error)}")
                    
                    # If no results from either search, try simple fallback
                    if not all_results:
                        print("🔄 No results from main searches, trying simple fallback")
                        try:
                            fallback_results = await search_documents_simple(query, 20, "fallback")
                            for result in fallback_results:
                                all_results[result.id] = result
                            print(f"✅ Fallback search returned {len(fallback_results)} results")
                        except Exception as fallback_error:
                            print(f"❌ Fallback search also failed: {str(fallback_error)}")
                    
                    # Convert to list and limit to 20 total results
                    final_results = list(all_results.values())[:20]
                    
                    if not final_results:
                        # Return a helpful message if no results found
                        return {
                            "jsonrpc": "2.0",
                            "id": request_id,
                            "result": {
                                "content": [
                                    {
                                        "type": "text",
                                        "text": json.dumps({
                                            "results": [],
                                            "message": f"No documents found for query: '{query}'. Try different search terms or check if documents exist in the library."
                                        }, indent=2)
                                    }
                                ]
                            }
                        }
                    
                    # Convert to the expected format
                    results_data = []
                    for result in final_results:
                        results_data.append({
                            "id": result.id,
                            "title": result.title,
                            "text": result.text,
                            "url": result.url
                        })
                    
                    print(f"✅ Returning {len(results_data)} search results")
                    
                    return {
                        "jsonrpc": "2.0",
                        "id": request_id,
                        "result": {
                            "content": [
                                {
                                    "type": "text",
                                    "text": json.dumps({"results": results_data}, indent=2)
                                }
                            ]
                        }
                    }
                    
                except Exception as e:
                    print(f"❌ Search failed: {str(e)}")
                    return {
                        "jsonrpc": "2.0",
                        "id": request_id,
                        "error": {
                            "code": -32603,
                            "message": f"Search failed: {str(e)}"
                        }
                    }
            
            elif tool_name == "fetch":
                try:
                    doc_id = arguments.get("id", "")
                    if not doc_id:
                        raise ValueError("Document ID parameter is required")
                    
                    print(f"📥 Fetching document: {doc_id}")
                    
                    document = await fetch_document_content(doc_id)
                    
                    print(f"✅ Document fetched successfully")
                    
                    return {
                        "jsonrpc": "2.0",
                        "id": request_id,
                        "result": {
                            "content": [
                                {
                                    "type": "text",
                                    "text": json.dumps(document, indent=2)
                                }
                            ]
                        }
                    }
                    
                except Exception as e:
                    print(f"❌ Fetch failed: {str(e)}")
                    return {
                        "jsonrpc": "2.0",
                        "id": request_id,
                        "error": {
                            "code": -32603,
                            "message": f"Fetch failed: {str(e)}"
                        }
                    }
            
            else:
                print(f"❌ Unknown tool: {tool_name}")
                return {
                    "jsonrpc": "2.0",
                    "id": request_id,
                    "error": {
                        "code": -32601,
                        "message": f"Unknown tool: {tool_name}"
                    }
                }
        
        else:
            print(f"❌ Unknown method: {method}")
            return {
                "jsonrpc": "2.0",
                "id": request_id,
                "error": {
                    "code": -32601,
                    "message": f"Unknown method: {method}"
                }
            }
    
    except json.JSONDecodeError as e:
        print(f"❌ Invalid JSON: {str(e)}")
        return {
            "jsonrpc": "2.0",
            "error": {
                "code": -32700,
                "message": "Parse error: Invalid JSON"
            }
        }
    
    except Exception as e:
        print(f"❌ Handler error: {str(e)}")
        return {
            "jsonrpc": "2.0",
            "id": request_id,
            "error": {
                "code": -32603,
                "message": f"Internal error: {str(e)}"
            }
        }

@app.get("/")
async def root():
    """Health check and basic info endpoint for GET requests"""
    print("🏥 Health check requested (GET)")
    return {
        "name": "iManage Deep Research MCP Server",
        "version": "1.0.0",
        "description": "MCP server for ChatGPT integration with iManage Work API",
        "protocol": "MCP/1.0",
        "capabilities": ["tools"],
        "status": "healthy",
        "authentication": "none",
        "endpoints": {
            "mcp": "POST /",
            "health": "GET /health",
            "test": "GET /test"
        }
    }

# ---- Manifest and Discovery Endpoints ----
@app.get("/manifest.json")
async def manifest():
    """Application manifest for ChatGPT"""
    print("📋 Manifest requested")
    return {
        "schema_version": "v1",
        "name_for_model": "imanage_deep_research",
        "name_for_human": "iManage Deep Research",
        "description_for_model": "Search and retrieve documents from iManage Work for deep research analysis",
        "description_for_human": "Access your iManage documents for comprehensive research",
        "auth": {
            "type": "none"
        },
        "api": {
            "type": "mcp",
            "url": "/",
            "has_user_authentication": False
        },
        "logo_url": "",
        "contact_email": "",
        "legal_info_url": ""
    }

@app.get("/.well-known/ai-plugin.json")
async def ai_plugin():
    """AI Plugin manifest"""
    print("🤖 AI Plugin manifest requested")
    return {
        "schema_version": "v1",
        "name_for_model": "imanage_deep_research",
        "name_for_human": "iManage Deep Research",
        "description_for_model": "Search and retrieve documents from iManage Work for deep research analysis. No authentication required.",
        "description_for_human": "Access your iManage documents for comprehensive research",
        "auth": {
            "type": "none"
        },
        "api": {
            "type": "mcp",
            "url": "/",
            "has_user_authentication": False
        }
    }
# ---- MCP Discovery Endpoints ----
@app.get("/.well-known/mcp")
async def mcp_discovery():
    """MCP discovery endpoint"""
    print("🔍 MCP discovery requested")
    return {
        "version": "1.0.0",
        "name": "iManage Deep Research MCP Server",
        "description": "Deep research connector for iManage Work API",
        "capabilities": {
            "tools": True,
            "resources": False,
            "prompts": False
        },
        "authentication": {
            "type": "none"
        },
        "endpoint": {
            "url": "/",
            "method": "POST"
        }
    }

@app.options("/")
async def options_handler():
    """Handle CORS preflight requests"""
    print("🔄 CORS preflight request")
    return {
        "Allow": "GET, POST, OPTIONS",
        "Access-Control-Allow-Origin": "*",
        "Access-Control-Allow-Methods": "GET, POST, OPTIONS",
        "Access-Control-Allow-Headers": "Content-Type, Authorization"
    }

# ---- OAuth/Connector Endpoints ----
@app.get("/oauth/authorize")
async def oauth_authorize():
    """Handle OAuth authorization - auto-approve since no real auth needed"""
    print("🔓 OAuth authorize request - auto-approving")
    return {
        "code": "auto_approved",
        "state": "no_auth_required"
    }

@app.post("/oauth/token") 
async def oauth_token():
    """Handle OAuth token request - return dummy token"""
    print("🔓 OAuth token request - returning dummy token")
    return {
        "access_token": "no_auth_required",
        "token_type": "bearer",
        "expires_in": 3600,
        "scope": "read"
    }

@app.get("/oauth/userinfo")
async def oauth_userinfo():
    """Handle OAuth user info request"""
    print("🔓 OAuth userinfo request")
    return {
        "sub": "imanage_user",
        "name": "iManage User",
        "email": "user@imanage.com"
    }

# ---- ChatGPT Connector Specific Endpoints ----
@app.get("/connectors/oauth")
async def connectors_oauth(request: Request):
    """Handle ChatGPT connector OAuth flow"""
    print("🔗 ChatGPT connector OAuth request")
    base_url = str(request.base_url).rstrip('/')
    return {
        "authorization_url": f"{base_url}/oauth/authorize",
        "token_url": f"{base_url}/oauth/token",
        "userinfo_url": f"{base_url}/oauth/userinfo",
        "scopes": [],
        "auto_approved": True
    }

@app.post("/connectors/oauth")
async def connectors_oauth_post():
    """Handle ChatGPT connector OAuth POST"""
    print("🔗 ChatGPT connector OAuth POST request")
    return {
        "access_token": "no_auth_required",
        "token_type": "bearer",
        "expires_in": 3600
    }
# ---- Authentication Endpoints ----
@app.get("/auth/status")
async def auth_status():
    """Return authentication status - no auth required"""
    print("🔓 Authentication status requested")
    return {
        "authenticated": True,
        "type": "none",
        "message": "No authentication required - server handles iManage auth internally"
    }

@app.post("/auth/callback")
async def auth_callback():
    """Handle auth callback - always successful since no auth needed"""
    print("✅ Auth callback - auto-success (no auth required)")
    return {
        "success": True,
        "authenticated": True,
        "message": "Authentication successful"
    }

# ---- Handle Any Auth-Related Requests ----
@app.get("/auth/{path:path}")
async def auth_catch_all(path: str):
    """Catch all auth requests and return success"""
    print(f"🔓 Auth catch-all request: {path}")
    return {
        "authenticated": True,
        "status": "success",
        "message": "No authentication required"
    }

@app.post("/auth/{path:path}")
async def auth_catch_all_post(path: str):
    """Catch all auth POST requests and return success"""
    print(f"🔓 Auth catch-all POST request: {path}")
    return {
        "authenticated": True,
        "status": "success",
        "message": "No authentication required"
    }

# ---- Legacy Endpoints for Testing ----
@app.get("/mcp/tools")
async def get_tools_legacy():
    """Legacy endpoint for testing - redirects to proper MCP format"""
    print("🔄 Legacy tools endpoint accessed - redirecting to MCP format")
    return {
        "message": "This is a legacy endpoint. Use POST / with proper MCP protocol.",
        "mcp_format": "POST / with method='tools/list'",
        "tools": [
            {
                "name": "search",
                "description": "Search iManage documents"
            },
            {
                "name": "fetch", 
                "description": "Fetch document content"
            }
        ]
    }

@app.get("/health")
async def health_check():
    """Simple health check endpoint"""
    print("🩺 Health check via /health")
    return {"status": "healthy", "timestamp": time.time()}

# ---- Test Endpoints ----
@app.get("/test")
async def test_connection():
    """Test iManage connection"""
    print("🧪 Testing iManage connection...")
    
    try:
        token = await get_token()
        
        # Test a simple API call
        test_url = f"{URL_PREFIX}/api/v2/customers/{CUSTOMER_ID}/features"
        headers = {"X-Auth-Token": token}
        
        async with httpx.AsyncClient(timeout=30.0) as client:
            response = await client.get(test_url, headers=headers)
            response.raise_for_status()
            
        print("✅ iManage connection test successful")
        return {
            "status": "success",
            "message": "iManage connection working",
            "customer_id": CUSTOMER_ID,
            "library_id": LIBRARY_ID,
            "timestamp": time.time()
        }
        
    except Exception as e:
        print(f"❌ iManage connection test failed: {str(e)}")
        return {
            "status": "error",
            "message": f"iManage connection failed: {str(e)}",
            "timestamp": time.time()
        }

@app.get("/test/search")
async def test_search():
    """Test a basic search to verify API format"""
    print("🧪 Testing basic search...")
    
    try:
        token = await get_token()
        
        # Test both URL formats
        search_urls = [
            f"{URL_PREFIX}/api/v2/customers/{CUSTOMER_ID}/libraries/{LIBRARY_ID}/documents",
            f"{URL_PREFIX}/work/web/api/v2/customers/{CUSTOMER_ID}/libraries/{LIBRARY_ID}/documents"
        ]
        
        results = {}
        headers = {"X-Auth-Token": token}
        
        async with httpx.AsyncClient(timeout=30.0) as client:
            for url in search_urls:
                try:
                    response = await client.get(url, headers=headers, params={"limit": 5})
                    results[url] = {
                        "status_code": response.status_code,
                        "response_sample": response.text[:500],
                        "success": response.status_code == 200
                    }
                except Exception as e:
                    results[url] = {
                        "error": str(e),
                        "success": False
                    }
            
            return {
                "status": "completed",
                "test_results": results,
                "recommendation": "Use the URL format that returned 200 status"
            }
            
    except Exception as e:
        print(f"❌ Search test failed: {str(e)}")
        return {
            "status": "error",
            "message": f"Search test failed: {str(e)}",
            "timestamp": time.time()
        }

@app.get("/test/document/{doc_id}")
async def test_document_access(doc_id: str):
    """Test document access with both URL formats"""
    print(f"🧪 Testing document access for: {doc_id}")
    
    try:
        token = await get_token()
        
        # Test both URL formats for document access
        doc_urls = [
            f"{URL_PREFIX}/api/v2/customers/{CUSTOMER_ID}/libraries/{LIBRARY_ID}/documents/{doc_id}",
            f"{URL_PREFIX}/work/web/api/v2/customers/{CUSTOMER_ID}/libraries/{LIBRARY_ID}/documents/{doc_id}"
        ]
        
        results = {}
        headers = {"X-Auth-Token": token}
        
        async with httpx.AsyncClient(timeout=30.0) as client:
            for url in doc_urls:
                try:
                    response = await client.get(url, headers=headers)
                    results[url] = {
                        "status_code": response.status_code,
                        "response_sample": response.text[:300],
                        "success": response.status_code == 200,
                        "accessible": "data" in response.text.lower()
                    }
                except Exception as e:
                    results[url] = {
                        "error": str(e),
                        "success": False,
                        "accessible": False
                    }
            
            return {
                "status": "completed",
                "document_id": doc_id,
                "test_results": results,
                "recommendation": "Use the URL format that returned 200 status and contains 'data'"
            }
            
    except Exception as e:
        @app.get("/test/processing")
        async def test_document_processing():
            """Test document processing capabilities"""
            print("🧪 Testing document processing capabilities...")
            
            capabilities = {
                "libraries_available": {
                    "PyPDF2": False,
                    "pdfplumber": False,
                    "python-docx": False,
                    "openpyxl": False,
                    "python-pptx": False,
                    "beautifulsoup4": False
                },
                "supported_formats": [],
                "overall_status": DOCUMENT_PROCESSING_AVAILABLE
            }
    
    # Test each library
    try:
        import PyPDF2
        capabilities["libraries_available"]["PyPDF2"] = True
        capabilities["supported_formats"].append("PDF")
    except ImportError:
        pass
    
    try:
        import pdfplumber
        capabilities["libraries_available"]["pdfplumber"] = True
    except ImportError:
        pass
    
    try:
        from docx import Document
        capabilities["libraries_available"]["python-docx"] = True
        capabilities["supported_formats"].append("Word Documents (.docx)")
    except ImportError:
        pass
    
    try:
        import openpyxl
        capabilities["libraries_available"]["openpyxl"] = True
        capabilities["supported_formats"].append("Excel Spreadsheets (.xlsx)")
    except ImportError:
        pass
    
    try:
        from pptx import Presentation
        capabilities["libraries_available"]["python-pptx"] = True
        capabilities["supported_formats"].append("PowerPoint Presentations (.pptx)")
    except ImportError:
        pass
    
    try:
        from bs4 import BeautifulSoup
        capabilities["libraries_available"]["beautifulsoup4"] = True
        capabilities["supported_formats"].append("HTML Documents")
    except ImportError:
        pass
    
    capabilities["supported_formats"].append("Plain Text")
    capabilities["supported_formats"].append("JSON")
    capabilities["supported_formats"].append("XML")
    
    return {
        "status": "completed",
        "document_processing": capabilities,
        "recommendation": "Install missing libraries with: pip install PyPDF2 pdfplumber python-docx openpyxl python-pptx beautifulsoup4" if not DOCUMENT_PROCESSING_AVAILABLE else "All document processing libraries are available"
    }

# ---- Startup ----
@app.on_event("startup")
async def startup_event():
    """Server startup logging"""
    print("🎉 iManage Deep Research MCP Server starting up")
    print(f"📁 Connected to Customer: {CUSTOMER_ID}, Library: {LIBRARY_ID}")
    
    # Test authentication on startup
    try:
        await get_token()
        print("✅ Initial authentication test successful")
    except Exception as e:
        print(f"⚠️ Warning: Initial authentication test failed: {str(e)}")

if __name__ == "__main__":
    import uvicorn
    
    print("🚀 Starting iManage Deep Research MCP Server...")
    uvicorn.run(
        app, 
        host="0.0.0.0", 
        port=int(os.getenv("PORT", 8000)),
        log_level="info"
    )